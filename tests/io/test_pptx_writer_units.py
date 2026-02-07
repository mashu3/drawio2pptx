"""
Unit tests for pptx_writer: create_presentation, add_slide with Shape/Connector/Text elements.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor

from drawio2pptx.io.pptx_writer import PPTXWriter
from drawio2pptx.model.intermediate import (
    ShapeElement,
    ConnectorElement,
    TextElement,
    Style,
    Transform,
    TextParagraph,
    TextRun,
)


def test_create_presentation_none_page_size() -> None:
    writer = PPTXWriter()
    prs, layout = writer.create_presentation(None)
    assert prs is not None
    assert layout is not None


def test_create_presentation_with_page_size() -> None:
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    assert prs is not None
    assert prs.slide_width is not None
    assert prs.slide_height is not None


def test_create_presentation_zero_page_size() -> None:
    """Zero page size: should not crash; slide size may stay default."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((0.0, 0.0))
    assert prs is not None


def test_add_slide_shape_element() -> None:
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="s1",
        x=50.0,
        y=50.0,
        w=100.0,
        h=60.0,
        shape_type="rectangle",
        style=Style(fill="default"),
    )
    writer.add_slide(prs, layout, [shape])
    assert len(prs.slides) == 1
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_shape_zero_size_not_added() -> None:
    """Shape with w<=0 or h<=0 is not added (early return in _add_shape)."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="s1",
        x=50.0,
        y=50.0,
        w=0.0,
        h=60.0,
        shape_type="rectangle",
        style=Style(fill="default"),
    )
    writer.add_slide(prs, layout, [shape])
    # No shape should be added (zero width)
    slide = prs.slides[0]
    names = [s.name for s in slide.shapes if getattr(s, "name", "").startswith("drawio2pptx:shape:")]
    assert len(names) == 0


def test_add_slide_line_shape() -> None:
    """Line shape uses _add_line_shape path."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="line1",
        x=10.0,
        y=10.0,
        w=200.0,
        h=2.0,
        shape_type="line",
        style=Style(stroke=RGBColor(0, 0, 0)),
    )
    writer.add_slide(prs, layout, [shape])
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_rounded_rectangle() -> None:
    """Shape with corner_radius uses ROUNDED_RECTANGLE."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="r1",
        x=50.0,
        y=50.0,
        w=100.0,
        h=60.0,
        shape_type="rectangle",
        style=Style(fill="default", corner_radius=10.0),
    )
    writer.add_slide(prs, layout, [shape])
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_connector_element() -> None:
    """Connector with points is added as connector (freeform or line)."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    conn = ConnectorElement(
        id="c1",
        source_id="s0",
        target_id="s1",
        points=[(100.0, 100.0), (200.0, 150.0)],
        edge_style="straight",
        style=Style(stroke=RGBColor(0, 0, 0)),
    )
    writer.add_slide(prs, layout, [conn])
    slide = prs.slides[0]
    # Connector may be emitted as group or line; we just check no crash and slide has content
    assert len(slide.shapes) >= 1


def test_add_slide_text_element() -> None:
    """Standalone TextElement is added via _add_text."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    text_el = TextElement(
        id="t1",
        x=50.0,
        y=50.0,
        w=200.0,
        h=30.0,
        text=[TextParagraph(runs=[TextRun(text="Hello")])],
    )
    writer.add_slide(prs, layout, [text_el])
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_shape_with_fill_rgb() -> None:
    """Shape with RGB fill applies solid fill."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="s1",
        x=50.0,
        y=50.0,
        w=80.0,
        h=50.0,
        shape_type="rectangle",
        style=Style(fill=RGBColor(255, 0, 0)),
    )
    writer.add_slide(prs, layout, [shape])
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_shape_no_stroke() -> None:
    """Shape with no_stroke disables line."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="s1",
        x=50.0,
        y=50.0,
        w=80.0,
        h=50.0,
        shape_type="rectangle",
        style=Style(fill="default", no_stroke=True),
    )
    writer.add_slide(prs, layout, [shape])
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_shape_has_shadow() -> None:
    """Shape with has_shadow enables shadow."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="s1",
        x=50.0,
        y=50.0,
        w=80.0,
        h=50.0,
        shape_type="rectangle",
        style=Style(fill="default", has_shadow=True),
    )
    writer.add_slide(prs, layout, [shape])
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_parallelogram() -> None:
    """Parallelogram shape gets adjustment set."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shape = ShapeElement(
        id="p1",
        x=50.0,
        y=50.0,
        w=100.0,
        h=50.0,
        shape_type="parallelogram",
        style=Style(fill="default"),
    )
    writer.add_slide(prs, layout, [shape])
    assert len(prs.slides[0].shapes) >= 1


def test_add_slide_multiple_elements() -> None:
    """Slide with multiple elements adds all in order."""
    writer = PPTXWriter()
    prs, layout = writer.create_presentation((800.0, 600.0))
    shapes = [
        ShapeElement(id="a", x=10.0, y=10.0, w=50.0, h=30.0, shape_type="rectangle", style=Style(fill="default")),
        ShapeElement(id="b", x=70.0, y=10.0, w=50.0, h=30.0, shape_type="ellipse", style=Style(fill="default")),
    ]
    writer.add_slide(prs, layout, shapes)
    assert len(prs.slides[0].shapes) >= 2
