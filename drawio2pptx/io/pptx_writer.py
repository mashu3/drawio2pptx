"""
PowerPoint output module

Generates PowerPoint presentations from intermediate models using python-pptx + lxml
"""
from typing import List, Optional, Tuple
from lxml import etree as ET
from pptx import Presentation  # type: ignore[import]
from pptx.util import Emu, Pt  # type: ignore[import]
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR  # type: ignore[import]
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore[import]
from pptx.dml.color import RGBColor  # type: ignore[import]
import io

from ..model.intermediate import BaseElement, ShapeElement, ConnectorElement, TextElement, TextParagraph, TextRun, ImageData
from ..media.image_utils import get_image_size, pad_image_to_square, prepare_image_for_pptx
from ..geom.units import px_to_emu, px_to_pt, scale_font_size_for_pptx
from ..geom.transform import split_polyline_to_segments
from ..mapping.shape_map import map_shape_type_to_pptx
from ..mapping.style_map import map_arrow_type, map_arrow_type_with_size, map_dash_pattern
from ..logger import ConversionLogger
from ..fonts import replace_font, DRAWIO_DEFAULT_FONT_FAMILY
from ..config import PARALLELOGRAM_SKEW, SWIMLANE_DEFAULT_PADDING_PX, ConversionConfig, default_config

# XML namespaces
NS_DRAWINGML = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_PRESENTATIONML = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NSMAP_DRAWINGML = {'a': NS_DRAWINGML}
NSMAP_PRESENTATIONML = {'p': NS_PRESENTATIONML}
NSMAP_BOTH = {'p': NS_PRESENTATIONML, 'a': NS_DRAWINGML}


def _a(tag_name: str) -> str:
    """Create DrawingML namespace-qualified tag name"""
    return f'{{{NS_DRAWINGML}}}{tag_name}'


def _p(tag_name: str) -> str:
    """Create PresentationML namespace-qualified tag name"""
    return f'{{{NS_PRESENTATIONML}}}{tag_name}'


class PPTXWriter:
    """PowerPoint presentation writer"""
    
    def __init__(self, logger: Optional[ConversionLogger] = None, config: Optional[ConversionConfig] = None):
        """
        Args:
            logger: ConversionLogger instance
            config: ConversionConfig instance (uses default_config if None)
        """
        self.config = config or default_config
        self.logger = logger
        self._svg_backend_logged = False

    def _set_shape_name(self, shape_obj, name: Optional[str]) -> None:
        """Set debug name on a shape/connector/textbox; log on failure."""
        if not name:
            return
        try:
            shape_obj.name = name
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set shape name: {e}")

    def _safe_try(self, fn, debug_msg: str) -> None:
        """Run fn(); on Exception log debug_msg and continue."""
        try:
            fn()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to {debug_msg}: {e}")

    @staticmethod
    def _shape_type_is(shape: ShapeElement, *names: str) -> bool:
        """Return True if shape.shape_type (normalized lower) equals any of *names."""
        return (shape.shape_type or "").strip().lower() in names

    @staticmethod
    def _shape_type_contains(shape: ShapeElement, substr: str) -> bool:
        """Return True if shape.shape_type (lower) contains substr."""
        return substr.lower() in (shape.shape_type or "").lower()

    def create_presentation(self, page_size: Optional[Tuple[float, float]] = None) -> Tuple[Presentation, object]:
        """
        Create presentation and blank layout.

        Args:
            page_size: (width, height) tuple (px), or None

        Returns:
            Tuple of (Presentation, blank layout).
        """
        prs = Presentation()
        
        # Get blank layout
        blank_layout_index = 6
        try:
            blank_layout = prs.slide_layouts[blank_layout_index]
        except Exception:
            blank_layout = prs.slide_layouts[0]
        
        # Set slide size
        if page_size and page_size[0] and page_size[1]:
            try:
                prs.slide_width = px_to_emu(page_size[0])
                prs.slide_height = px_to_emu(page_size[1])
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set slide size: {e}")
        
        return prs, blank_layout

    def add_slide(self, prs: Presentation, blank_layout, elements: List[BaseElement]):
        """
        Add elements to slide
        
        Args:
            prs: Presentation object
            blank_layout: Blank layout
            elements: List of elements (sorted by Z-order; later elements are on top)
        """
        slide = prs.slides.add_slide(blank_layout)
        
        # Add elements in the provided stacking order (later = topmost in PowerPoint).
        for element in elements:
            if isinstance(element, ShapeElement):
                self._add_shape(slide, element)
            elif isinstance(element, ConnectorElement):
                self._add_connector(slide, element)
            elif isinstance(element, TextElement):
                self._add_text(slide, element)
    
    def _compute_shape_geometry(self, shape: ShapeElement) -> Tuple[int, int, int, int]:
        """Compute (left_emu, top_emu, width_emu, height_emu) for add_shape, including step/arrow adjustments."""
        left = px_to_emu(shape.x)
        top = px_to_emu(shape.y)
        width = px_to_emu(shape.w)
        height = px_to_emu(shape.h)
        try:
            if self._shape_type_is(shape, "step"):
                step_size = getattr(shape.style, "step_size_px", None)
                if step_size is not None and shape.w > 0:
                    gap_px = max(0.0, min(step_size * 0.1, shape.w * 0.3))
                    width = px_to_emu(max(shape.w - gap_px, 1.0))
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to apply step gap: {e}")
        try:
            if self._shape_type_is(shape, "right_arrow", "notched_right_arrow"):
                rot = float(getattr(shape.transform, "rotation", 0.0) or 0.0) % 360.0
                if abs(rot - 90.0) < 1.0 or abs(rot - 270.0) < 1.0:
                    width, height = height, width
                    left = px_to_emu(shape.x + (shape.w - shape.h) / 2.0)
                    top = px_to_emu(shape.y + (shape.h - shape.w) / 2.0)
        except Exception:
            pass
        return left, top, width, height

    def _add_shape(self, slide, shape: ShapeElement):
        """Add shape"""
        if shape.w <= 0 or shape.h <= 0:
            return None

        if self._shape_type_is(shape, "line"):
            return self._add_line_shape(slide, shape)

        pptx_shape_type = map_shape_type_to_pptx(shape.shape_type)
        if (pptx_shape_type == MSO_SHAPE.RECTANGLE and
                shape.style.corner_radius is not None and
                shape.style.corner_radius > 0):
            pptx_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

        left, top, width, height = self._compute_shape_geometry(shape)
        shp = slide.shapes.add_shape(
            pptx_shape_type,
            left, top, width, height
        )
        
        self._set_shape_name(shp, f"drawio2pptx:shape:{shape.id}" if shape.id else None)
        self._apply_shape_adjustments(shp, shape)

        # If the shape is flipped (flipH/flipV), flipping via <a:xfrm> also flips the text in PowerPoint.
        # draw.io typically keeps label glyphs upright. Work around by rendering the label as a separate
        # transparent textbox overlay (not flipped), while keeping the shape geometry flipped.
        shape_is_flipped = False
        try:
            shape_is_flipped = bool(getattr(shape.transform, "flip_h", False) or getattr(shape.transform, "flip_v", False))
        except Exception:
            shape_is_flipped = False

        # When verticalLabelPosition=bottom, label is rendered below the shape (see _add_shape_label_below).
        # For image shapes, text should always be below the image (y-axis direction), not inside the shape.
        # Do not set text inside the shape in that case.
        label_below = (
            shape.text
            and not shape_is_flipped
            and (
                getattr(shape.style, "vertical_label_position", None) == "bottom"
                or shape.image  # For image shapes, always place text below the image
            )
        )

        # Set text (inside the shape when not flipped and not label-below)
        if shape.text and not shape_is_flipped and not label_below:
            margin_overrides, text_direction = self._get_swimlane_text_options(shape)
            if margin_overrides is None and bool(getattr(shape.style, "aws_group_text_padding", False)):
                # AWS group containers tend to have top-aligned labels; ensure small inner padding
                # so text does not stick to the border.
                first_para = shape.text[0] if shape.text else None
                top_px = float(getattr(first_para, "spacing_top", 0.0) or 0.0)
                left_px = float(getattr(first_para, "spacing_left", 0.0) or 0.0)
                bottom_px = float(getattr(first_para, "spacing_bottom", 0.0) or 0.0)
                right_px = float(getattr(first_para, "spacing_right", 0.0) or 0.0)
                margin_overrides = (
                    max(top_px, 6.0),
                    max(left_px, 8.0),
                    max(bottom_px, 2.0),
                    max(right_px, 2.0),
                )
            self._set_text_frame(
                shp.text_frame,
                shape.text,
                default_highlight_color=shape.style.label_background_color,
                word_wrap=shape.style.word_wrap,
                margin_overrides_px=margin_overrides,
                text_direction=text_direction,
                clip_overflow=bool(getattr(shape.style, "clip_text", False)),
                shape_height_px=shape.h,
            )
        
        # If shape has an image, make the shape completely invisible (no fill, no stroke)
        # The image will be added as a separate picture shape below
        if shape.image:
            # Make shape completely transparent (no fill)
            try:
                shp.fill.background()
            except Exception:
                pass
            # Remove stroke completely (no border)
            try:
                self._set_no_line_xml(shp)
            except Exception:
                try:
                    shp.line.fill.background()
                except Exception:
                    pass
            stroke_color = None
        else:
            self._apply_shape_fill(shp, shape)
            stroke_color = self._apply_shape_stroke(shp, shape)
        
        self._maybe_add_swimlane_divider(slide, shape, stroke_color)
        self._apply_shape_shadow(shp, shape.style.has_shadow)

        self._maybe_add_cube_3d(shp, shape)
        self._maybe_add_bpmn_symbol(slide, shape)
        
        # Add image if present (as separate picture shape)
        if shape.image:
            self._safe_try(lambda: self._add_shape_image(slide, shape), "add shape image")
        else:
            self._safe_try(
                lambda: self._add_aws_group_icon_overlay(slide, shape),
                "add aws group icon overlay",
            )
        
        # Add label below image if verticalLabelPosition=bottom or if shape has an image
        # For image shapes, text is always placed below the image (y-axis direction)
        if label_below:
            self._safe_try(lambda: self._add_shape_label_below(slide, shape), "add shape label below")
        if shape.text and shape_is_flipped:
            self._safe_try(lambda: self._add_shape_text_overlay(slide, shape), "add flipped-shape text overlay")
        return shp

    def _add_aws_group_icon_overlay(self, slide, shape: ShapeElement):
        """Add a small top-left overlay icon for aws4 group/groupCenter containers."""
        icon_ref = getattr(shape.style, "aws_group_icon_ref", None)
        if not icon_ref:
            return None

        # AWS group icons are square badges attached to container border.
        icon_size_px = max(14.0, min(24.0, min(float(shape.w), float(shape.h)) * 0.18))
        icon_key = (getattr(shape.style, "aws_group_icon_key", None) or "").lower()
        if icon_key == "group_auto_scaling_group":
            # Special case: auto scaling icon sits on top edge, horizontally centered.
            left_px = shape.x + max((shape.w - icon_size_px) / 2.0, 0.0)
            top_px = shape.y
        else:
            left_px = shape.x
            top_px = shape.y

        left = px_to_emu(left_px)
        top = px_to_emu(top_px)
        width = px_to_emu(icon_size_px)
        height = px_to_emu(icon_size_px)

        data_uri = icon_ref if icon_ref.startswith("data:") else None
        file_path = None if data_uri else icon_ref

        image_bytes, img_width_px, img_height_px, _ = prepare_image_for_pptx(
            data_uri=data_uri,
            file_path=file_path,
            shape_type=shape.shape_type,
            target_width_px=int(icon_size_px),
            target_height_px=int(icon_size_px),
            base_dpi=self.config.dpi if hasattr(self.config, "dpi") else 192.0,
            aws_icon_color_hex=None,
        )
        if not image_bytes:
            return None

        # Keep original icon aspect and pad to square instead of stretching.
        padding_color_hex = None
        stroke = getattr(shape.style, "stroke", None)
        padding_color_mode = (getattr(shape.style, "aws_group_icon_padding_color_mode", None) or "stroke").lower()
        if isinstance(stroke, RGBColor) and padding_color_mode != "icon":
            padding_color_hex = f"{stroke[0]:02X}{stroke[1]:02X}{stroke[2]:02X}"
        padding_ratio = getattr(shape.style, "aws_group_icon_padding_ratio", None)
        try:
            padding_ratio = float(padding_ratio) if padding_ratio is not None else 0.18
        except (TypeError, ValueError):
            padding_ratio = 0.18
        image_bytes = pad_image_to_square(
            image_bytes,
            padding_ratio=padding_ratio,
            padding_color_hex=padding_color_hex,
        )
        sq_w, sq_h = get_image_size(image_bytes)
        if not sq_w or not sq_h:
            return None

        picture = slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
        self._set_shape_name(
            picture,
            f"drawio2pptx:aws-group-icon:{shape.id}" if shape.id else None,
        )
        try:
            picture.line.fill.background()
        except Exception:
            pass
        return picture

    def _add_transparent_textbox(self, slide, left_emu, top_emu, width_emu, height_emu, name: Optional[str] = None):
        """Add a textbox with transparent fill and line; return the shape."""
        tb = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
        self._set_shape_name(tb, name)
        try:
            tb.fill.background()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set textbox fill background: {e}")
        try:
            tb.line.fill.background()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set textbox line background: {e}")
        return tb

    def _add_shape_text_overlay(self, slide, shape: ShapeElement):
        """Render shape label as a transparent textbox overlay (used for flipped shapes)."""
        if not shape.text or shape.w <= 0 or shape.h <= 0:
            return None
        left = px_to_emu(shape.x)
        top = px_to_emu(shape.y)
        width = px_to_emu(shape.w)
        height = px_to_emu(shape.h)
        tb = self._add_transparent_textbox(
            slide, left, top, width, height,
            f"drawio2pptx:shape-text-overlay:{shape.id}" if shape.id else None,
        )

        # Apply the same paragraph formatting (reuse swimlane margin/direction when applicable).
        margin_overrides, text_direction = self._get_swimlane_text_options(shape)
        self._set_text_frame(
            tb.text_frame,
            shape.text,
            default_highlight_color=shape.style.label_background_color,
            word_wrap=shape.style.word_wrap,
            margin_overrides_px=margin_overrides,
            text_direction=text_direction,
            clip_overflow=bool(getattr(shape.style, "clip_text", False)),
            shape_height_px=shape.h,
        )

        return tb

    def _add_shape_label_below(self, slide, shape: ShapeElement):
        """Render shape label in a textbox below the shape (draw.io verticalLabelPosition=bottom)."""
        if not shape.text or shape.w <= 0 or shape.h <= 0:
            return None

        gap_px = 4.0
        label_y = shape.y + shape.h + gap_px

        # Estimate label height from actual text content and font size
        # Count lines across all paragraphs
        lines: List[str] = []
        font_size_px = 12.0
        for para in shape.text:
            text = "".join(run.text for run in para.runs if run.text)
            if text:
                lines.extend(text.splitlines() or [text])
            # Get font size from first run found
            if font_size_px == 12.0:  # Still default
                for run in para.runs:
                    if run.font_size is not None:
                        font_size_px = float(run.font_size)
                        break

        if not lines:
            lines = [""]

        # Calculate height based on number of lines and font size
        # Use similar logic to _estimate_text_box_size: line_height = font_size * 1.4
        label_h_px = max(len(lines) * float(font_size_px) * 1.4, float(font_size_px) * 1.2)

        left = px_to_emu(shape.x)
        top = px_to_emu(label_y)
        width = px_to_emu(shape.w)
        height = px_to_emu(label_h_px)
        tb = self._add_transparent_textbox(
            slide, left, top, width, height,
            f"drawio2pptx:shape-label-below:{shape.id}" if shape.id else None,
        )

        self._set_text_frame(
            tb.text_frame,
            shape.text,
            default_highlight_color=shape.style.label_background_color,
            word_wrap=False,
            margin_overrides_px=None,
            text_direction=None,
            clip_overflow=False,
        )
        return tb

    def _apply_shape_fill(self, shp, shape: ShapeElement) -> None:
        """Apply fill (swimlane gradient, default, solid, or background) to a shape."""
        fill_color = shape.style.fill
        is_swimlane = getattr(shape.style, "is_swimlane", False)
        swimlane_start = float(getattr(shape.style, "swimlane_start_size", 0) or 0)
        if is_swimlane and swimlane_start > 0 and shape.h > 0 and isinstance(fill_color, RGBColor):
            self._set_swimlane_gradient_fill_xml(shp, shape)
        elif fill_color == "default":
            self._set_default_fill_xml(shp)
        elif fill_color:
            try:
                shp.fill.solid()
                shp.fill.fore_color.rgb = fill_color
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set fill color: {e}")
        else:
            try:
                shp.fill.background()
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set background fill: {e}")

    def _maybe_add_swimlane_divider(
        self, slide, shape: ShapeElement, stroke_color: Optional[RGBColor]
    ) -> None:
        """Add swimlane header divider line when applicable."""
        try:
            if not getattr(shape.style, "is_swimlane", False):
                return
            # Don't draw divider line when swimlane has strokeColor=none (to avoid drawing non-existent black line)
            if stroke_color is None:
                return
            start_size = float(getattr(shape.style, "swimlane_start_size", 0.0) or 0.0)
            if start_size <= 0 or not getattr(shape.style, "swimlane_line", True):
                return
            self._add_swimlane_header_divider(
                slide=slide,
                shape=shape,
                stroke_color=stroke_color,
                stroke_width_px=shape.style.stroke_width,
            )
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to add swimlane header divider: {e}")

    def _apply_shape_stroke(self, shp, shape: ShapeElement) -> Optional[RGBColor]:
        """Apply stroke (or no_stroke) to a shape. Returns stroke color for use by e.g. swimlane divider, or None if no stroke."""
        if getattr(shape.style, "no_stroke", False):
            try:
                self._set_no_line_xml(shp)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to disable stroke: {e}")
            return None
        stroke_color = shape.style.stroke if shape.style.stroke else RGBColor(0, 0, 0)
        try:
            shp.line.fill.solid()
            self._set_stroke_color_xml(shp, stroke_color)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set stroke color: {e}")
        if shape.style.stroke_width > 0:
            try:
                shp.line.width = px_to_pt(shape.style.stroke_width)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set stroke width: {e}")
        # Apply dash pattern for vertex shapes (e.g. dashed containers/boundaries)
        dash_pattern = getattr(shape.style, "dash", None)
        if dash_pattern:
            try:
                self._set_dash_pattern_xml(shp, dash_pattern)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set shape dash pattern: {e}")
        return stroke_color

    def _apply_shape_shadow(self, shp, has_shadow: bool) -> None:
        """Enable or disable shadow on a shape (inherit from theme or disable via XML)."""
        try:
            if has_shadow:
                shp.shadow.inherit = True
            else:
                shp.shadow.inherit = False
                self._disable_shadow_xml(shp)
        except Exception:
            if not has_shadow:
                self._disable_shadow_xml(shp)

    def _get_swimlane_text_options(self, shape: ShapeElement) -> Tuple[Optional[Tuple[float, float, float, float]], Optional[str]]:
        """Return (margin_overrides_px, text_direction) for swimlane shapes; otherwise (None, None)."""
        try:
            if not getattr(shape.style, "is_swimlane", False):
                return None, None
            start_size = float(getattr(shape.style, "swimlane_start_size", 0.0) or 0.0)
            if start_size <= 0 or not shape.text:
                return None, None
            first_para = shape.text[0]
            base_top = first_para.spacing_top if first_para.spacing_top is not None else 0.0
            base_left = first_para.spacing_left if first_para.spacing_left is not None else 0.0
            base_bottom = first_para.spacing_bottom if first_para.spacing_bottom is not None else 0.0
            base_right = first_para.spacing_right if first_para.spacing_right is not None else 0.0
            # Apply default inner padding for header when draw.io did not specify spacing (box looks flush otherwise)
            pad = SWIMLANE_DEFAULT_PADDING_PX
            if base_top == 0 and base_left == 0 and base_right == 0:
                base_top = pad
                base_left = pad
                base_right = pad
            if getattr(shape.style, "swimlane_horizontal", False):
                margin_overrides = (
                    base_top,
                    base_left,
                    max(shape.h - start_size + base_bottom, 0.0),
                    base_right,
                )
                return margin_overrides, None
            text_direction = "vert270"
            margin_overrides = (
                base_top,
                base_left,
                base_bottom,
                max(shape.w - start_size + base_right, 0.0),
            )
            return margin_overrides, text_direction
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to compute swimlane text options: {e}")
            return None, None

    def _apply_shape_transform(self, shp, shape: ShapeElement) -> None:
        """Apply shape rotation/flip from the intermediate model to the PPTX shape."""
        if not hasattr(shape, "transform") or shape.transform is None:
            return

        # Rotation (degrees). python-pptx supports shape.rotation for AutoShapes.
        try:
            rot = float(getattr(shape.transform, "rotation", 0.0) or 0.0)
            if abs(rot) > 1e-6:
                try:
                    shp.rotation = rot
                except Exception:
                    pass
        except Exception:
            pass

        # Flip (OpenXML: a:xfrm@flipH / a:xfrm@flipV). No public python-pptx API; set via XML.
        flip_h = bool(getattr(shape.transform, "flip_h", False))
        flip_v = bool(getattr(shape.transform, "flip_v", False))
        if flip_h or flip_v:
            self._set_flip_xml(shp, flip_h=flip_h, flip_v=flip_v)

    def _set_flip_xml(self, shape, flip_h: bool = False, flip_v: bool = False) -> None:
        """Set flipH/flipV on <a:xfrm> via XML for preset shapes."""
        try:
            if not hasattr(shape, "_element"):
                return
            shape_element = shape._element

            xfrm = shape_element.find(".//a:xfrm", namespaces=NSMAP_DRAWINGML)
            if xfrm is None:
                # Fallback: namespace-qualified tag search
                xfrm = shape_element.find(f".//{_a('xfrm')}")
            if xfrm is None:
                return

            if flip_h:
                xfrm.set("flipH", "1")
            else:
                xfrm.attrib.pop("flipH", None)

            if flip_v:
                xfrm.set("flipV", "1")
            else:
                xfrm.attrib.pop("flipV", None)
        except Exception:
            return None

    def _maybe_add_cube_3d(self, shp, shape: ShapeElement) -> None:
        """Apply cube 3D rotation XML when shape type is cube."""
        if not self._shape_type_is(shape, "cube"):
            return
        self._safe_try(lambda: self._set_cube_3d_rotation_xml(shp), "set cube 3D rotation")

    def _maybe_add_bpmn_symbol(self, slide, shape: ShapeElement) -> None:
        """Add BPMN symbol overlay (e.g. parallel gateway) when style requests it."""
        bpmn = getattr(shape.style, "bpmn_symbol", None)
        if not bpmn or (bpmn or "").strip().lower() != "parallelgw":
            return
        self._safe_try(lambda: self._add_bpmn_parallel_gateway_symbol(slide, shape), "add BPMN symbol overlay")

    def _apply_shape_adjustments(self, shp, shape: ShapeElement) -> None:
        """Apply parallelogram skew, step chevron size, and rotation/flip to a shape."""
        try:
            if self._shape_type_contains(shape, "parallelogram") and hasattr(shp, "adjustments") and len(shp.adjustments) > 0:
                shp.adjustments[0] = float(PARALLELOGRAM_SKEW)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set parallelogram adjustments: {e}")
        try:
            if self._shape_type_is(shape, "step") and hasattr(shp, "adjustments") and len(shp.adjustments) > 0:
                step_size = getattr(shape.style, "step_size_px", None)
                if step_size is not None and shape.w > 0:
                    adj = max(0.02, min(0.6, (step_size / float(shape.w)) * 1.5))
                    shp.adjustments[0] = float(adj)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set step chevron adjustment: {e}")
        try:
            self._apply_shape_transform(shp, shape)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to apply shape transform: {e}")

    def _apply_line_style_to_connector(
        self,
        line_shape,
        stroke_color: RGBColor,
        stroke_width_px: float,
        has_shadow: bool,
        log_prefix: str = "line",
    ) -> None:
        """Apply stroke color, width, and shadow to a connector/line shape."""
        try:
            line_shape.line.fill.solid()
            self._set_stroke_color_xml(line_shape, stroke_color)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set {log_prefix} stroke color: {e}")
        if stroke_width_px and stroke_width_px > 0:
            try:
                line_shape.line.width = px_to_pt(stroke_width_px)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set {log_prefix} width: {e}")
        try:
            if has_shadow:
                line_shape.shadow.inherit = True
            else:
                line_shape.shadow.inherit = False
                self._disable_shadow_xml(line_shape)
        except Exception:
            if not has_shadow:
                self._disable_shadow_xml(line_shape)

    def _add_line_shape(self, slide, shape: ShapeElement):
        """Add line shape (draw.io line vertex) as a connector."""
        if shape.w >= shape.h:
            y = shape.y + (shape.h / 2.0)
            x1, y1 = shape.x, y
            x2, y2 = shape.x + shape.w, y
        else:
            x = shape.x + (shape.w / 2.0)
            x1, y1 = x, shape.y
            x2, y2 = x, shape.y + shape.h

        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1), px_to_emu(y1), px_to_emu(x2), px_to_emu(y2),
        )
        self._set_shape_name(line, f"drawio2pptx:line:{shape.id}" if shape.id else None)
        stroke_color = shape.style.stroke if shape.style.stroke else RGBColor(0, 0, 0)
        self._apply_line_style_to_connector(
            line, stroke_color, shape.style.stroke_width, shape.style.has_shadow, log_prefix="line",
        )
        return line

    def _add_swimlane_header_divider(
        self,
        slide,
        shape: ShapeElement,
        stroke_color: RGBColor,
        stroke_width_px: float,
    ):
        """Draw the swimlane header divider line."""
        start_size = float(getattr(shape.style, "swimlane_start_size", 0.0) or 0.0)
        if start_size <= 0:
            return None

        is_horizontal = bool(getattr(shape.style, "swimlane_horizontal", False))
        if is_horizontal:
            # Header on top: horizontal divider at y + startSize
            x1 = shape.x
            y1 = shape.y + start_size
            x2 = shape.x + shape.w
            y2 = y1
        else:
            # Header on left: vertical divider at x + startSize
            x1 = shape.x + start_size
            y1 = shape.y
            x2 = x1
            y2 = shape.y + shape.h

        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1),
            px_to_emu(y1),
            px_to_emu(x2),
            px_to_emu(y2),
        )
        self._set_shape_name(line, f"drawio2pptx:swimlane-divider:{shape.id}" if shape.id else None)
        has_shadow = getattr(shape.style, "has_shadow", False)
        self._apply_line_style_to_connector(line, stroke_color, stroke_width_px, has_shadow, log_prefix="swimlane divider")
        return line
    
    def _add_connector(self, slide, connector: ConnectorElement):
        """Add connector as a single polyline shape."""
        if not connector.points or len(connector.points) < 2:
            return None

        # Orthogonal connectors should preserve segment structure to match draw.io stacking
        # and for reliable z-order testing/debugging (each segment gets its own name).
        try:
            if (connector.edge_style or "").lower() == "orthogonal":
                return self._add_orthogonal_connector(slide, connector)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to add orthogonal connector: {e}")

        # Simplify nearly straight polylines to avoid jitter in Freeform.
        points_px = self._simplify_polyline_points(connector.points, tol_px=0.5)
        if len(points_px) >= 2 and self._is_almost_straight(points_px, tol_px=0.5):
            return self._add_straight_connector(slide, connector, points_px)

        (
            line_points_px,
            add_open_oval_start,
            add_open_oval_end,
            effective_start_arrow,
            effective_end_arrow,
            start_marker_center_px,
            end_marker_center_px,
        ) = self._compute_open_oval_trimmed_points(connector, points_px)
        start_arrow = connector.style.arrow_start
        end_arrow = connector.style.arrow_end

        # For normal lines (straight), use FreeformBuilder
        # Convert points to EMU
        points_emu = [(px_to_emu(x), px_to_emu(y)) for x, y in line_points_px]
        
        try:
            if len(points_emu) < 2:
                return None
            
            # Create polyline with FreeformBuilder
            x0, y0 = points_emu[0]
            builder = slide.shapes.build_freeform(x0, y0)
            
            if len(points_emu) > 1:
                remaining_points = points_emu[1:]
                builder.add_line_segments(remaining_points, close=False)
            
            line_shape = builder.convert_to_shape()
        except Exception:
            return None

        self._set_shape_name(line_shape, f"drawio2pptx:connector:{connector.id}" if connector.id else None)

        self._apply_connector_line_style(
            line_shape, connector,
            effective_start_arrow=effective_start_arrow,
            effective_end_arrow=effective_end_arrow,
        )

        self._add_open_oval_markers_for_connector(
            slide, connector,
            add_open_oval_start, add_open_oval_end,
            start_marker_center_px, end_marker_center_px,
            start_arrow, end_arrow,
        )

        return line_shape

    def _compute_open_oval_trimmed_points(
        self, connector: ConnectorElement, points_px: List[Tuple[float, float]]
    ) -> Tuple[
        List[Tuple[float, float]], bool, bool, Optional[str], Optional[str],
        Tuple[float, float], Tuple[float, float],
    ]:
        """Compute trimmed line points and open-oval flags. Returns (line_points_px, add_start, add_end, eff_start_arrow, eff_end_arrow, start_center_px, end_center_px)."""
        start_arrow = connector.style.arrow_start
        end_arrow = connector.style.arrow_end
        add_open_oval_start = self._should_emulate_open_oval_marker(start_arrow, connector.style.arrow_start_fill)
        add_open_oval_end = self._should_emulate_open_oval_marker(end_arrow, connector.style.arrow_end_fill)
        effective_start_arrow = None if add_open_oval_start else start_arrow
        effective_end_arrow = None if add_open_oval_end else end_arrow
        start_center_px = points_px[0]
        end_center_px = points_px[-1]
        line_points_px = list(points_px)
        if add_open_oval_start or add_open_oval_end:
            try:
                start_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_start_size_px,
                    )
                    if add_open_oval_start
                    else 0.0
                )
                end_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_end_size_px,
                    )
                    if add_open_oval_end
                    else 0.0
                )
                line_points_px = self._trim_polyline_endpoints_px(line_points_px, start_trim, end_trim)
            except Exception:
                line_points_px = list(points_px)
        return (
            line_points_px,
            add_open_oval_start,
            add_open_oval_end,
            effective_start_arrow,
            effective_end_arrow,
            start_center_px,
            end_center_px,
        )

    def _add_straight_connector(self, slide, connector: ConnectorElement, points_px: List[Tuple[float, float]]):
        """Add a single straight connector (line) between endpoints."""
        if len(points_px) < 2:
            return None

        (
            line_points_px,
            add_open_oval_start,
            add_open_oval_end,
            effective_start_arrow,
            effective_end_arrow,
            start_marker_center_px,
            end_marker_center_px,
        ) = self._compute_open_oval_trimmed_points(connector, points_px)
        start_arrow = connector.style.arrow_start
        end_arrow = connector.style.arrow_end

        if len(line_points_px) < 2:
            return None
        (x1, y1), (x2, y2) = line_points_px[0], line_points_px[-1]
        line_shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(x1),
            px_to_emu(y1),
            px_to_emu(x2),
            px_to_emu(y2),
        )

        self._set_shape_name(line_shape, f"drawio2pptx:connector:{connector.id}" if connector.id else None)

        self._apply_connector_line_style(
            line_shape, connector,
            effective_start_arrow=effective_start_arrow,
            effective_end_arrow=effective_end_arrow,
        )

        self._add_open_oval_markers_for_connector(
            slide, connector,
            add_open_oval_start, add_open_oval_end,
            start_marker_center_px, end_marker_center_px,
            start_arrow, end_arrow,
        )

        return line_shape

    def _add_open_oval_markers_for_connector(
        self,
        slide,
        connector: ConnectorElement,
        add_open_oval_start: bool,
        add_open_oval_end: bool,
        start_center_px: Tuple[float, float],
        end_center_px: Tuple[float, float],
        start_arrow: Optional[str],
        end_arrow: Optional[str],
    ) -> None:
        """Add open-oval overlay markers at connector endpoints when needed."""
        marker_configs = [
            ("start", add_open_oval_start, start_center_px, start_arrow, connector.style.arrow_start_size_px),
            ("end", add_open_oval_end, end_center_px, end_arrow, connector.style.arrow_end_size_px),
        ]
        for position, should_add, center_px, arrow_name, arrow_size_px in marker_configs:
            if should_add:
                x, y = center_px
                self._add_open_oval_marker(
                    slide=slide,
                    x_px=x,
                    y_px=y,
                    stroke_color=connector.style.stroke,
                    stroke_width_px=connector.style.stroke_width,
                    arrow_name=arrow_name,
                    arrow_size_px=arrow_size_px,
                    marker_name=f"drawio2pptx:marker:open-oval:{connector.id}:{position}",
                )

    def _apply_connector_line_style(
        self,
        line_shape,
        connector: ConnectorElement,
        effective_start_arrow: Optional[str] = None,
        effective_end_arrow: Optional[str] = None,
        arrow_start_fill: Optional[bool] = None,
        arrow_end_fill: Optional[bool] = None,
        arrow_start_size_px: Optional[float] = None,
        arrow_end_size_px: Optional[float] = None,
    ) -> None:
        """Apply stroke, fill off, dash, arrows, and shadow to a connector/line shape.
        Arrow fill/size default to connector.style when omitted (e.g. for single-segment connectors).
        """
        style = connector.style
        stroke_color = style.stroke if style.stroke else RGBColor(0, 0, 0)
        try:
            line_shape.line.fill.solid()
            self._set_edge_stroke_color_xml(line_shape, stroke_color)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set connector stroke color: {e}")
        if style.stroke_width > 0:
            try:
                line_shape.line.width = px_to_pt(style.stroke_width)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set connector stroke width: {e}")
        try:
            line_shape.fill.background()
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to disable connector fill: {e}")
        if style.dash:
            try:
                self._set_dash_pattern_xml(line_shape, style.dash)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set connector dash pattern: {e}")
        start_fill = arrow_start_fill if arrow_start_fill is not None else style.arrow_start_fill
        end_fill = arrow_end_fill if arrow_end_fill is not None else style.arrow_end_fill
        start_sz = arrow_start_size_px if arrow_start_size_px is not None else style.arrow_start_size_px
        end_sz = arrow_end_size_px if arrow_end_size_px is not None else style.arrow_end_size_px
        # When endSize/startSize omitted, use larger arrow for thick lines (e.g. strokeWidth 3)
        if style.stroke_width >= 2.5:
            if start_sz is None and effective_start_arrow:
                start_sz = 12.0  # maps to "lg"
            if end_sz is None and effective_end_arrow:
                end_sz = 12.0  # maps to "lg"
        if effective_start_arrow or effective_end_arrow:
            self._set_arrow_heads_xml(
                line_shape, effective_start_arrow, effective_end_arrow,
                start_fill, end_fill, style.stroke, start_sz, end_sz,
            )
        try:
            if style.has_shadow:
                line_shape.shadow.inherit = True
            else:
                line_shape.shadow.inherit = False
                self._disable_shadow_xml(line_shape)
        except Exception:
            if not style.has_shadow:
                self._disable_shadow_xml(line_shape)

    @staticmethod
    def _simplify_polyline_points(points: List[Tuple[float, float]], tol_px: float = 0.5) -> List[Tuple[float, float]]:
        """Drop nearly collinear points to stabilize Freeform rendering."""
        if len(points) <= 2:
            return list(points)

        def _dist_point_to_line(px, py, ax, ay, bx, by):
            dx = bx - ax
            dy = by - ay
            denom = dx * dx + dy * dy
            if denom <= 1e-9:
                return ((px - ax) ** 2 + (py - ay) ** 2) ** 0.5
            t = ((px - ax) * dx + (py - ay) * dy) / denom
            t = 0.0 if t < 0.0 else 1.0 if t > 1.0 else t
            cx = ax + t * dx
            cy = ay + t * dy
            return ((px - cx) ** 2 + (py - cy) ** 2) ** 0.5

        simplified = [points[0]]
        for i in range(1, len(points) - 1):
            ax, ay = simplified[-1]
            bx, by = points[i + 1]
            px, py = points[i]
            if _dist_point_to_line(px, py, ax, ay, bx, by) > tol_px:
                simplified.append(points[i])
        simplified.append(points[-1])
        return simplified

    @staticmethod
    def _is_almost_straight(points: List[Tuple[float, float]], tol_px: float = 0.5) -> bool:
        """Return True if all points are close to the line from start to end."""
        if len(points) <= 2:
            return True
        (ax, ay) = points[0]
        (bx, by) = points[-1]
        dx = bx - ax
        dy = by - ay
        denom = dx * dx + dy * dy
        if denom <= 1e-9:
            return True
        for (px, py) in points[1:-1]:
            t = ((px - ax) * dx + (py - ay) * dy) / denom
            cx = ax + t * dx
            cy = ay + t * dy
            dist = ((px - cx) ** 2 + (py - cy) ** 2) ** 0.5
            if dist > tol_px:
                return False
        return True

    def _add_text(self, slide, text_element: TextElement):
        """Add standalone text element."""
        if text_element.w <= 0 or text_element.h <= 0:
            return None

        left = px_to_emu(text_element.x)
        top = px_to_emu(text_element.y)
        width = px_to_emu(text_element.w)
        height = px_to_emu(text_element.h)
        tb = self._add_transparent_textbox(
            slide, left, top, width, height,
            f"drawio2pptx:text:{text_element.id}" if text_element.id else None,
        )

        if text_element.text:
            self._set_text_frame(
                tb.text_frame,
                text_element.text,
                default_highlight_color=text_element.style.label_background_color,
                word_wrap=text_element.style.word_wrap,
                clip_overflow=bool(getattr(text_element.style, "clip_text", False)),
            )

        return tb

    def _add_bpmn_parallel_gateway_symbol(self, slide, shape: ShapeElement):
        """Add plus sign overlay for BPMN parallel gateway using lines (not text)."""
        symbol_size = min(shape.w, shape.h) * 1.0
        center_x = shape.x + shape.w / 2.0
        center_y = shape.y + shape.h / 2.0

        stroke_color = shape.style.stroke if shape.style.stroke else RGBColor(0, 0, 0)
        line_width_pt = 2.0
        half_length = symbol_size * 0.25

        h_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(center_x - half_length),
            px_to_emu(center_y),
            px_to_emu(center_x + half_length),
            px_to_emu(center_y),
        )
        v_line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px_to_emu(center_x),
            px_to_emu(center_y - half_length),
            px_to_emu(center_x),
            px_to_emu(center_y + half_length),
        )

        try:
            h_line.line.fill.solid()
            self._set_edge_stroke_color_xml(h_line, stroke_color)
            h_line.line.width = Pt(line_width_pt)
            self._remove_arrowheads_xml(h_line)
            h_line.shadow.inherit = False
            self._disable_shadow_xml(h_line)
            self._remove_effect_ref_xml(h_line)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set horizontal line properties: {e}")

        self._set_shape_name(h_line, f"drawio2pptx:bpmn-symbol-h:{shape.id}" if shape.id else None)

        try:
            v_line.line.fill.solid()
            self._set_edge_stroke_color_xml(v_line, stroke_color)
            v_line.line.width = Pt(line_width_pt)
            self._remove_arrowheads_xml(v_line)
            v_line.shadow.inherit = False
            self._disable_shadow_xml(v_line)
            self._remove_effect_ref_xml(v_line)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set vertical line properties: {e}")

        self._set_shape_name(v_line, f"drawio2pptx:bpmn-symbol-v:{shape.id}" if shape.id else None)
    
    def _add_orthogonal_connector(self, slide, connector: ConnectorElement):
        """Add polyline as straight connectors for each segment"""
        if not connector.points or len(connector.points) < 2:
            return None

        # Open-oval marker emulation needs endpoint trimming on the underlying polyline.
        add_open_oval_start = self._should_emulate_open_oval_marker(connector.style.arrow_start, connector.style.arrow_start_fill)
        add_open_oval_end = self._should_emulate_open_oval_marker(connector.style.arrow_end, connector.style.arrow_end_fill)

        start_marker_center_px = connector.points[0]
        end_marker_center_px = connector.points[-1]
        points_for_segments = list(connector.points)
        if add_open_oval_start or add_open_oval_end:
            try:
                start_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_start_size_px,
                    )
                    if add_open_oval_start
                    else 0.0
                )
                end_trim = (
                    self._open_oval_trim_radius_px(
                        stroke_width_px=connector.style.stroke_width,
                        arrow_size_px=connector.style.arrow_end_size_px,
                    )
                    if add_open_oval_end
                    else 0.0
                )
                points_for_segments = self._trim_polyline_endpoints_px(points_for_segments, start_trim, end_trim)
            except Exception:
                points_for_segments = list(connector.points)
        
        # Split polyline into segments
        segments = split_polyline_to_segments(points_for_segments)
        if not segments:
            return None
        
        created_shapes = []
        
        for idx, ((x1, y1), (x2, y2)) in enumerate(segments):
            try:
                x1_emu, y1_emu = px_to_emu(x1), px_to_emu(y1)
                x2_emu, y2_emu = px_to_emu(x2), px_to_emu(y2)
                
                # Create each segment as a straight connector
                conn_shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    x1_emu, y1_emu, x2_emu, y2_emu
                )

                self._set_shape_name(conn_shape, f"drawio2pptx:connector:{connector.id}:seg:{idx}" if connector.id else None)

                is_first_segment = idx == 0
                is_last_segment = idx == len(segments) - 1
                start_arrow = connector.style.arrow_start if is_first_segment else None
                end_arrow = connector.style.arrow_end if is_last_segment else None
                if start_arrow or end_arrow:
                    add_open_oval_start = self._should_emulate_open_oval_marker(
                        start_arrow, connector.style.arrow_start_fill if is_first_segment else False
                    )
                    add_open_oval_end = self._should_emulate_open_oval_marker(
                        end_arrow, connector.style.arrow_end_fill if is_last_segment else False
                    )
                    eff_start = None if add_open_oval_start else start_arrow
                    eff_end = None if add_open_oval_end else end_arrow
                else:
                    eff_start, eff_end = None, None

                self._apply_connector_line_style(
                    conn_shape, connector,
                    effective_start_arrow=eff_start,
                    effective_end_arrow=eff_end,
                    arrow_start_fill=connector.style.arrow_start_fill if is_first_segment else False,
                    arrow_end_fill=connector.style.arrow_end_fill if is_last_segment else False,
                    arrow_start_size_px=connector.style.arrow_start_size_px if is_first_segment else None,
                    arrow_end_size_px=connector.style.arrow_end_size_px if is_last_segment else None,
                )

                created_shapes.append(conn_shape)
            except Exception as e:
                if self.logger:
                    self.logger.warning(f"Failed to create connector segment {idx}: {e}")
                continue
        
        if not created_shapes:
            return None

        # Add open-oval markers once, after all segments are created (ensures the marker is on top
        # of the line geometry for this connector).
        marker_configs = {
            "start": (add_open_oval_start, start_marker_center_px, connector.style.arrow_start, connector.style.arrow_start_size_px),
            "end": (add_open_oval_end, end_marker_center_px, connector.style.arrow_end, connector.style.arrow_end_size_px),
        }
        for position, (should_add, center_px, arrow_name, arrow_size_px) in marker_configs.items():
            if should_add:
                x, y = center_px
                self._add_open_oval_marker(
                    slide=slide,
                    x_px=x,
                    y_px=y,
                    stroke_color=connector.style.stroke,
                    stroke_width_px=connector.style.stroke_width,
                    arrow_name=arrow_name,
                    arrow_size_px=arrow_size_px,
                    marker_name=f"drawio2pptx:marker:open-oval:{connector.id}:{position}",
                )
        
        return created_shapes[0]

    @staticmethod
    def _should_emulate_open_oval_marker(arrow_name: Optional[str], fill: bool) -> bool:
        """Return True when we should emulate an open oval marker with an overlay shape.

        PowerPoint's line-end (a:headEnd/a:tailEnd) doesn't support an "unfilled oval" variant.
        draw.io represents it with startArrow/endArrow=oval and startFill/endFill=0.
        """
        if not arrow_name:
            return False
        try:
            return arrow_name.strip().lower() == "oval" and (fill is False)
        except Exception:
            return False

    def _add_open_oval_marker(
        self,
        slide,
        x_px: float,
        y_px: float,
        stroke_color: Optional[RGBColor],
        stroke_width_px: float,
        arrow_name: Optional[str],
        arrow_size_px: Optional[float],
        marker_name: str,
    ):
        """Add a small ellipse outline at the given point to emulate an unfilled oval line-end.

        The connector geometry is trimmed so the line stops at the marker boundary, so we keep the marker
        unfilled (noFill) rather than using a white "mask" fill (which breaks on non-white backgrounds).
        """
        try:
            # Prefer draw.io marker size (startSize/endSize) when present.
            # In mxGraph style, marker size values are expressed in screen px.
            # For an oval marker, treating the size as an approximate diameter matches draw.io more closely.
            #
            # Important: draw.io can omit startSize/endSize from the style string. In that case,
            # mxGraph defaults the marker size to 6. Use that default here to avoid oversizing.
            try:
                effective_size_px = float(arrow_size_px) if arrow_size_px is not None else 6.0
            except Exception:
                effective_size_px = 6.0

            base_d = max(effective_size_px, 1.0)
            d_px = max(base_d, 6.0 + float(stroke_width_px) * 1.25)

            left = px_to_emu(x_px - d_px / 2.0)
            top = px_to_emu(y_px - d_px / 2.0)
            width = px_to_emu(d_px)
            height = px_to_emu(d_px)

            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            try:
                marker.name = marker_name
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker name: {e}")

            # No fill (hollow marker).
            try:
                marker.fill.background()
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker fill background: {e}")
            try:
                self._set_no_fill_xml(marker)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker no fill XML: {e}")

            # Stroke matches connector.
            # Use black as default if stroke_color is None
            effective_stroke_color = stroke_color if stroke_color is not None else RGBColor(0, 0, 0)
            try:
                marker.line.fill.solid()
                self._set_stroke_color_xml(marker, effective_stroke_color)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker stroke color: {e}")

            try:
                if stroke_width_px and stroke_width_px > 0:
                    marker.line.width = px_to_pt(stroke_width_px)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set marker stroke width: {e}")

            # Markers should not add extra visual effects.
            try:
                marker.shadow.inherit = False
                self._disable_shadow_xml(marker)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to disable marker shadow: {e}")
        except Exception:
            return None

    @staticmethod
    def _open_oval_marker_diameter_px(stroke_width_px: float, arrow_size_px: Optional[float]) -> float:
        """Return the diameter (px) for the emulated open-oval marker."""
        try:
            effective_size_px = float(arrow_size_px) if arrow_size_px is not None else 6.0
        except Exception:
            effective_size_px = 6.0
        base_d = max(effective_size_px, 1.0)
        try:
            return max(base_d, 6.0 + float(stroke_width_px) * 1.25)
        except Exception:
            return max(base_d, 6.0)

    @classmethod
    def _open_oval_trim_radius_px(cls, stroke_width_px: float, arrow_size_px: Optional[float]) -> float:
        """Return trim distance (px) so the connector line stops at the open-oval boundary."""
        d_px = cls._open_oval_marker_diameter_px(stroke_width_px=stroke_width_px, arrow_size_px=arrow_size_px)
        try:
            return max(d_px / 2.0 - float(stroke_width_px) / 2.0, 0.0)
        except Exception:
            return max(d_px / 2.0, 0.0)

    @staticmethod
    def _trim_polyline_endpoints_px(
        points: List[Tuple[float, float]],
        start_trim_px: float,
        end_trim_px: float,
    ) -> List[Tuple[float, float]]:
        """Trim polyline endpoints by the requested distances (px)."""
        if not points or len(points) < 2:
            return points

        def dist(a: Tuple[float, float], b: Tuple[float, float]) -> float:
            dx = b[0] - a[0]
            dy = b[1] - a[1]
            return (dx * dx + dy * dy) ** 0.5

        def lerp(a: Tuple[float, float], b: Tuple[float, float], t: float) -> Tuple[float, float]:
            return (a[0] + (b[0] - a[0]) * t, a[1] + (b[1] - a[1]) * t)

        pts = list(points)

        # Trim start
        remaining = max(float(start_trim_px or 0.0), 0.0)
        while remaining > 1e-6 and len(pts) >= 2:
            seg_len = dist(pts[0], pts[1])
            if seg_len < 1e-6:
                pts.pop(0)
                continue
            if seg_len <= remaining and len(pts) > 2:
                remaining -= seg_len
                pts.pop(0)
                continue
            t = min(remaining / seg_len, 0.9)
            pts[0] = lerp(pts[0], pts[1], t)
            break

        # Trim end
        remaining = max(float(end_trim_px or 0.0), 0.0)
        while remaining > 1e-6 and len(pts) >= 2:
            seg_len = dist(pts[-1], pts[-2])
            if seg_len < 1e-6:
                pts.pop(-1)
                continue
            if seg_len <= remaining and len(pts) > 2:
                remaining -= seg_len
                pts.pop(-1)
                continue
            t = min(remaining / seg_len, 0.9)
            pts[-1] = lerp(pts[-1], pts[-2], t)
            break

        if len(pts) < 2:
            return list(points)
        return pts

    def _set_no_fill_xml(self, shape):
        """Force <a:noFill/> on the shape fill (spPr) via XML."""
        try:
            if not hasattr(shape, "_element"):
                return

            shape_element = shape._element

            # spPr is often a direct child in python-pptx shapes.
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break

            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in sp_pr.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        sp_pr.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove fill element {tag}: {e}")

            ET.SubElement(sp_pr, _a("noFill"))
        except Exception:
            return None

    def _set_no_line_xml(self, shape):
        """Force <a:noFill/> on the line (ln) via XML."""
        try:
            if not hasattr(shape, "_element"):
                return

            shape_element = shape._element
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break
            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            ln_element = sp_pr.find(".//a:ln", namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                ln_element = ET.SubElement(sp_pr, _a("ln"))

            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in ln_element.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        ln_element.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove line element {tag}: {e}")

            ET.SubElement(ln_element, _a("noFill"))
        except Exception:
            return None

    def _add_paragraphs_to_text_frame(
        self,
        text_frame,
        paragraphs: List[TextParagraph],
        default_highlight_color: Optional[RGBColor] = None,
    ) -> None:
        """Add paragraphs and runs to a text frame (after clear). Sets spacing, alignment, font, color, highlight."""
        for para_data in paragraphs:
            p = text_frame.add_paragraph()
            try:
                p.space_before = Pt(float(para_data.space_before_pt)) if getattr(para_data, "space_before_pt", None) is not None else Pt(0)
            except Exception:
                p.space_before = Pt(0)
            try:
                p.space_after = Pt(float(para_data.space_after_pt)) if getattr(para_data, "space_after_pt", None) is not None else Pt(0)
            except Exception:
                p.space_after = Pt(0)
            try:
                p.line_spacing = float(para_data.line_spacing) if para_data.line_spacing is not None else 1.0
            except Exception:
                p.line_spacing = 1.0
            HORIZONTAL_ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
            p.alignment = HORIZONTAL_ALIGN_MAP.get(para_data.align or "center", PP_ALIGN.CENTER)
            for run_data in para_data.runs:
                run = p.add_run()
                run.text = run_data.text
                effective_font_family = run_data.font_family or DRAWIO_DEFAULT_FONT_FAMILY
                replaced_font = replace_font(effective_font_family, config=self.config)
                try:
                    run.font.name = replaced_font if replaced_font else effective_font_family
                except Exception as e:
                    if self.logger:
                        self.logger.warning(f"Failed to set font: {e}")
                if run_data.font_size:
                    run.font.size = Pt(scale_font_size_for_pptx(run_data.font_size))
                else:
                    run.font.size = Pt(scale_font_size_for_pptx(12.0))
                run.font.bold = run_data.bold
                run.font.italic = run_data.italic
                run.font.underline = run_data.underline
                if run_data.font_color:
                    self._set_font_color_xml(run, run_data.font_color)
                else:
                    self._set_font_color_xml(run, RGBColor(0, 0, 0))
                if default_highlight_color is not None:
                    self._set_highlight_color_xml(run, default_highlight_color)

    # Height (px) below which a shape is treated as a "row cell" (e.g. class diagram row).
    # For such cells, verticalAlign=top is overridden to middle so text is vertically centered.
    _ROW_CELL_HEIGHT_THRESHOLD_PX = 40.0

    def _set_text_frame(
        self,
        text_frame,
        paragraphs: List[TextParagraph],
        default_highlight_color: Optional[RGBColor] = None,
        word_wrap: bool = True,
        margin_overrides_px: Optional[Tuple[float, float, float, float]] = None,
        text_direction: Optional[str] = None,
        clip_overflow: bool = False,
        shape_height_px: Optional[float] = None,
    ):
        """Set text frame"""
        if not paragraphs:
            return
        
        # Configure text frame
        text_frame.word_wrap = word_wrap
        text_frame.auto_size = None
        
        # Optional text direction (used for vertical swimlane headers)
        if text_direction is not None:
            try:
                self._set_text_direction_xml(text_frame, text_direction)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set text direction: {e}")

        # Get padding from first paragraph
        first_para = paragraphs[0]
        if margin_overrides_px is not None:
            top_px, left_px, bottom_px, right_px = margin_overrides_px
        else:
            top_px = first_para.spacing_top or 0
            left_px = first_para.spacing_left or 0
            bottom_px = first_para.spacing_bottom or 0
            right_px = first_para.spacing_right or 0

        # When overflow=hidden (clip_overflow), ignore negative margins so text does not
        # get pushed outside the shape and appear to overflow at the top.
        if clip_overflow:
            try:
                top_px = max(float(top_px or 0), 0.0)
                left_px = max(float(left_px or 0), 0.0)
                bottom_px = max(float(bottom_px or 0), 0.0)
                right_px = max(float(right_px or 0), 0.0)
            except (TypeError, ValueError):
                top_px = max(top_px or 0, 0.0)
                left_px = max(left_px or 0, 0.0)
                bottom_px = max(bottom_px or 0, 0.0)
                right_px = max(right_px or 0, 0.0)

        text_frame.margin_top = px_to_emu(top_px or 0)
        text_frame.margin_left = px_to_emu(left_px or 0)
        text_frame.margin_bottom = px_to_emu(bottom_px or 0)
        text_frame.margin_right = px_to_emu(right_px or 0)
        margin_px = (top_px or 0, left_px or 0, bottom_px or 0, right_px or 0)
        
        # Set vertical anchor (similar to legacy: default is middle)
        VERTICAL_ALIGN_MAP = {
            "top": (MSO_ANCHOR.TOP, 't'),
            "middle": (MSO_ANCHOR.MIDDLE, 'ctr'),
            "bottom": (MSO_ANCHOR.BOTTOM, 'b'),
        }
        vertical_align = first_para.vertical_align or "middle"
        # For short row cells (e.g. class diagram rows), use middle so text is vertically centered
        if (
            shape_height_px is not None
            and shape_height_px <= self._ROW_CELL_HEIGHT_THRESHOLD_PX
            and vertical_align == "top"
        ):
            vertical_align = "middle"
        saved_vertical_anchor, anchor_value = VERTICAL_ALIGN_MAP.get(
            vertical_align, (MSO_ANCHOR.MIDDLE, 'ctr')
        )

        # Note: clip_overflow now uses normAutofit (shrink-to-fit) instead of hard-clipping,
        # so negative inset + anchor adjustments are no longer needed.
        
        # Clear existing paragraphs
        text_frame.clear()
        
        # Remove empty paragraphs (similar to legacy: to ensure vertical_anchor works correctly)
        while len(text_frame.paragraphs) > 0 and not text_frame.paragraphs[0].text:
            try:
                first_p = text_frame.paragraphs[0]
                first_p._element.getparent().remove(first_p._element)
            except Exception:
                break

        self._add_paragraphs_to_text_frame(text_frame, paragraphs, default_highlight_color)

        # Reset vertical_anchor (similar to legacy: set after clear())
        # vertical_anchor is reset after tf.clear(), so need to set again
        text_frame.vertical_anchor = saved_vertical_anchor
        
        # Set vertical anchor via XML (similar to legacy: workaround for python-pptx bug)
        self._set_vertical_anchor_xml(text_frame, anchor_value, word_wrap, margin_px, clip_overflow=clip_overflow)
        
        # Do not override per-paragraph spacing here; it is used to emulate HTML block margins.
    
    def _set_vertical_anchor_xml(
        self,
        text_frame,
        anchor_value: str,
        word_wrap: bool = True,
        margin_px: Optional[Tuple[float, float, float, float]] = None,
        clip_overflow: bool = False,
    ):
        """Set vertical anchor via XML"""
        try:
            body_pr = text_frame._element.find(f'.//{_a("bodyPr")}')
            if body_pr is not None:
                body_pr.set('anchor', anchor_value)
                if body_pr.get('anchorCtr') is not None:
                    body_pr.attrib.pop('anchorCtr', None)
                if margin_px is not None:
                    top_px, left_px, bottom_px, right_px = margin_px
                    body_pr.set('tIns', str(int(px_to_emu(top_px or 0))))
                    body_pr.set('lIns', str(int(px_to_emu(left_px or 0))))
                    body_pr.set('bIns', str(int(px_to_emu(bottom_px or 0))))
                    body_pr.set('rIns', str(int(px_to_emu(right_px or 0))))
                # Set wrap attribute based on word_wrap setting
                # 'square' enables wrapping, 'none' disables it
                wrap_value = 'square' if word_wrap else 'none'
                body_pr.set('wrap', wrap_value)
                # Overflow handling when requested (draw.io: overflow=hidden).
                # Instead of hard-clipping (which hides visible glyphs), use
                # <a:normAutofit/> to auto-shrink text so it fits inside the shape.
                # This keeps all text readable while respecting the bounding box.
                if clip_overflow:
                    body_pr.attrib.pop('vertOverflow', None)
                    body_pr.attrib.pop('horzOverflow', None)
                    # Remove any existing autofit children
                    for autofit_tag in ('noAutofit', 'normAutofit', 'spAutoFit'):
                        for existing in body_pr.findall(f'.//a:{autofit_tag}', namespaces=NSMAP_DRAWINGML):
                            try:
                                body_pr.remove(existing)
                            except Exception:
                                pass
                        for existing in body_pr.findall(f'{_a(autofit_tag)}'):
                            try:
                                body_pr.remove(existing)
                            except Exception:
                                pass
                    ET.SubElement(body_pr, _a('normAutofit'))
                else:
                    body_pr.attrib.pop('vertOverflow', None)
                    body_pr.attrib.pop('horzOverflow', None)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set vertical anchor XML: {e}")

    def _set_text_direction_xml(self, text_frame, direction: str):
        """Set text direction via XML (bodyPr@vert)."""
        if not direction:
            return
        try:
            body_pr = text_frame._element.find(f'.//{_a("bodyPr")}')
            if body_pr is None:
                return
            body_pr.set('vert', str(direction))
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set text direction XML: {e}")
    
    def _set_font_color_xml(self, run, font_color: RGBColor):
        """Set font color via XML (similar to legacy _set_font_color)"""
        if not font_color:
            return
        
        try:
            # First try normal method to set
            try:
                run.font.color.rgb = font_color
            except Exception:
                # If setting fails, we'll use XML method directly
                pass
            
            # Always use XML method to ensure color is set correctly
            # python-pptx's RGB setting doesn't always work reliably
            try:
                run_element = getattr(run, "_element", None) or getattr(run, "_r", None)
                if run_element is not None:
                    
                    # Find rPr (run properties) element
                    r_pr = run_element.find('.//a:rPr', namespaces=NSMAP_DRAWINGML)
                    if r_pr is None:
                        # Create rPr if it doesn't exist
                        r_pr = ET.SubElement(run_element, _a('rPr'))
                    
                    # Find or create solidFill element
                    solid_fill = r_pr.find('.//a:solidFill', namespaces=NSMAP_DRAWINGML)
                    if solid_fill is None:
                        solid_fill = ET.SubElement(r_pr, _a('solidFill'))
                    
                    # Remove existing color elements
                    for color_elem in solid_fill.findall('.//a:srgbClr', namespaces=NSMAP_DRAWINGML):
                        solid_fill.remove(color_elem)
                    for color_elem in solid_fill.findall('.//a:schemeClr', namespaces=NSMAP_DRAWINGML):
                        solid_fill.remove(color_elem)
                    
                    # Add srgbClr element
                    srgb = ET.SubElement(solid_fill, _a('srgbClr'))
                    val = f"{font_color[0]:02X}{font_color[1]:02X}{font_color[2]:02X}"
                    srgb.set('val', val)
            except Exception as e:
                if self.logger:
                    self.logger.debug(f"Failed to set font color XML: {e}")
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to access run element for font color: {e}")

    def _set_highlight_color_xml(self, run, highlight_color: RGBColor) -> None:
        """Set highlight color via XML"""
        if not highlight_color:
            return
        try:
            run_element = getattr(run, "_element", None) or getattr(run, "_r", None)
            if run_element is None:
                return

            r_pr = run_element.find(".//a:rPr", namespaces=NSMAP_DRAWINGML)
            if r_pr is None:
                r_pr = ET.SubElement(run_element, _a("rPr"))

            # Remove existing highlight elements (avoid duplicates).
            for hi in r_pr.findall(".//a:highlight", namespaces=NSMAP_DRAWINGML):
                try:
                    r_pr.remove(hi)
                except Exception as e:
                    if self.logger:
                        self.logger.debug(f"Failed to remove highlight element: {e}")

            # Important: DrawingML elements are order-sensitive in many Office apps.
            # Insert <a:highlight> before font elements (<a:latin>/<a:ea>/<a:cs>) when present,
            # otherwise after fill elements when present, otherwise append.
            hi = ET.Element(_a("highlight"))
            srgb = ET.SubElement(hi, _a("srgbClr"))
            val = f"{highlight_color[0]:02X}{highlight_color[1]:02X}{highlight_color[2]:02X}"
            srgb.set("val", val)

            children = list(r_pr)
            insert_idx = len(children)

            font_tags = {_a("latin"), _a("ea"), _a("cs"), _a("sym")}
            for i, child in enumerate(children):
                if child.tag in font_tags:
                    insert_idx = i
                    break

            if insert_idx == len(children):
                fill_tags = {_a("noFill"), _a("solidFill"), _a("gradFill"), _a("pattFill"), _a("blipFill")}
                last_fill = -1
                for i, child in enumerate(children):
                    if child.tag in fill_tags:
                        last_fill = i
                if last_fill != -1:
                    insert_idx = last_fill + 1

            r_pr.insert(insert_idx, hi)
        except Exception:
            return
    
    def _set_default_fill_xml(self, shape):
        """Set default fill (white)"""
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            # spPr exists as a direct child element
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith('}spPr') or 'spPr' in child.tag:
                    sp_pr = child
                    break
            
            if sp_pr is None:
                return
            
            # Remove existing fill elements (noFill, solidFill, etc.)
            for fill_elem in sp_pr.findall('.//a:noFill', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(fill_elem)
            for fill_elem in sp_pr.findall('.//a:solidFill', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(fill_elem)
            for fill_elem in sp_pr.findall('.//a:gradFill', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(fill_elem)
            
            # Create solidFill element
            solid_fill = ET.SubElement(sp_pr, _a('solidFill'))
            
            # Add srgbClr element with white color (RGB: 255, 255, 255)
            srgb_clr = ET.SubElement(solid_fill, _a('srgbClr'))
            srgb_clr.set('val', 'FFFFFF')
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set default fill XML: {e}")

    def _set_swimlane_gradient_fill_xml(self, shp, shape: ShapeElement) -> None:
        """
        Set swimlane fill: header = fillColor, body = swimlaneFillColor.
        Uses a multi-stop gradient so only the header area is colored; body stays white/transparent.
        """
        def _rgb_to_hex(rgb: RGBColor) -> str:
            return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

        try:
            if not hasattr(shp, "_element"):
                return
            shape_element = shp._element
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break
            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            start_size = float(getattr(shape.style, "swimlane_start_size", 0) or 0.0)
            is_horizontal = bool(getattr(shape.style, "swimlane_horizontal", True))
            header_color = shape.style.fill
            body_color = getattr(shape.style, "swimlane_fill_color", None)
            if isinstance(header_color, RGBColor):
                header_rgb = header_color
            else:
                header_rgb = RGBColor(0xE0, 0xE0, 0xE0)
            if isinstance(body_color, RGBColor):
                body_rgb = body_color
            elif body_color in ("default", "auto"):
                body_rgb = RGBColor(0xFF, 0xFF, 0xFF)
            else:
                body_rgb = RGBColor(0xFF, 0xFF, 0xFF)

            if is_horizontal and shape.h > 0:
                ratio = min(1.0, max(0.0, start_size / shape.h))
                pos_header = int(ratio * 100000)
                ang = "5400000"  # 90 degrees, top to bottom
            else:
                if shape.w > 0:
                    ratio = min(1.0, max(0.0, start_size / shape.w))
                else:
                    ratio = 0.0
                pos_header = int(ratio * 100000)
                ang = "0"  # 0 degrees, left to right

            # Exactly 2 stops on the divider line (e.g. 50%); distance between them almost 0.
            # Before first stop = header, between = tiny blend, after second = body.
            pos_header = min(max(0, pos_header), 100000)
            pos1 = max(0, pos_header - 1)
            pos2 = min(100000, pos_header + 1)

            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in sp_pr.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        sp_pr.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove fill element {tag}: {e}")

            grad_fill = ET.SubElement(sp_pr, _a("gradFill"))
            grad_fill.set("rotWithShape", "1")
            gs_lst = ET.SubElement(grad_fill, _a("gsLst"))

            # Stop 1: just above divider → header (extends to 0%)
            gs1 = ET.SubElement(gs_lst, _a("gs"))
            gs1.set("pos", str(pos1))
            ET.SubElement(gs1, _a("srgbClr")).set("val", _rgb_to_hex(header_rgb))
            # Stop 2: just below divider → body (extends to 100%); distance to pos1 ≈ 0
            gs2 = ET.SubElement(gs_lst, _a("gs"))
            gs2.set("pos", str(pos2))
            ET.SubElement(gs2, _a("srgbClr")).set("val", _rgb_to_hex(body_rgb))

            lin = ET.SubElement(grad_fill, _a("lin"))
            lin.set("ang", ang)
            lin.set("scaled", "1")
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set swimlane gradient fill: {e}")

    def _set_linear_gradient_fill_xml(
        self,
        shape,
        base_fill: Optional[object],
        gradient_color: Optional[object],
        gradient_direction: Optional[str] = None,
    ) -> None:
        """
        Set a simple 2-stop linear gradient via DrawingML XML (<a:gradFill>).

        Notes:
            - draw.io uses style keys: fillColor (start) and gradientColor (end).
            - When either side is "default", we use white (FFFFFF) for start, light gray (E0E0E0) for end.
            - This is best-effort; PowerPoint themes may render slightly differently than draw.io.
        """

        def _rgb_to_hex(rgb: RGBColor) -> str:
            return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

        def _darken(rgb: RGBColor, amount: float = 0.25) -> RGBColor:
            # amount in [0, 1]; 0.25 means 25% closer to black
            r = max(0, int(round(rgb[0] * (1.0 - amount))))
            g = max(0, int(round(rgb[1] * (1.0 - amount))))
            b = max(0, int(round(rgb[2] * (1.0 - amount))))
            return RGBColor(r, g, b)

        def _direction_to_ang(direction: Optional[str]) -> str:
            # OpenXML uses 60000ths of a degree.
            DIRECTION_DEGREE_MAP = {
                "east": 0,
                "right": 0,
                "west": 180,
                "left": 180,
                "north": 270,
                "up": 270,
                "south": 90,
                "down": 90,
            }
            d = (direction or "").strip().lower()
            deg = DIRECTION_DEGREE_MAP.get(d, 90)  # default: vertical (top->bottom)
            return str(int(deg * 60000))

        try:
            if not hasattr(shape, "_element"):
                return

            shape_element = shape._element

            # spPr is often a direct child in python-pptx shapes.
            sp_pr = None
            for child in shape_element:
                if child.tag.endswith("}spPr") or "spPr" in child.tag:
                    sp_pr = child
                    break
            if sp_pr is None:
                sp_pr = shape_element.find(".//a:spPr", namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return

            # Remove existing fill elements to avoid duplicates.
            for tag in ("noFill", "solidFill", "gradFill", "pattFill", "blipFill"):
                for elem in sp_pr.findall(f".//a:{tag}", namespaces=NSMAP_DRAWINGML):
                    try:
                        sp_pr.remove(elem)
                    except Exception as e:
                        if self.logger:
                            self.logger.debug(f"Failed to remove gradient fill element {tag}: {e}")

            grad_fill = ET.SubElement(sp_pr, _a("gradFill"))
            grad_fill.set("rotWithShape", "1")

            gs_lst = ET.SubElement(grad_fill, _a("gsLst"))

            # Determine start (fillColor) and end (gradientColor)
            start_is_default = (base_fill in (None, "default", "auto"))
            end_is_default = (gradient_color in ("default", "auto"))

            start_rgb = base_fill if isinstance(base_fill, RGBColor) else None
            end_rgb = gradient_color if isinstance(gradient_color, RGBColor) else None

            # If gradientColor is "default" but we have an explicit fillColor, derive a darker end.
            if end_is_default and start_rgb is not None:
                end_rgb = _darken(start_rgb, 0.25)
                end_is_default = False

            # Stop 1 (pos=0)
            gs1 = ET.SubElement(gs_lst, _a("gs"))
            gs1.set("pos", "0")
            if start_rgb is not None:
                clr1 = ET.SubElement(gs1, _a("srgbClr"))
                clr1.set("val", _rgb_to_hex(start_rgb))
            else:
                # Default: white
                clr1 = ET.SubElement(gs1, _a("srgbClr"))
                clr1.set("val", "FFFFFF")

            # Stop 2 (pos=100000)
            gs2 = ET.SubElement(gs_lst, _a("gs"))
            gs2.set("pos", "100000")
            if end_rgb is not None:
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", _rgb_to_hex(end_rgb))
            elif end_is_default or start_is_default:
                # Default: slightly darker gray (for gradient effect)
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", "E0E0E0")
            elif start_rgb is not None:
                # Fallback: derive darker from start
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", _rgb_to_hex(_darken(start_rgb, 0.25)))
            else:
                # Default: slightly darker gray (for gradient effect)
                clr2 = ET.SubElement(gs2, _a("srgbClr"))
                clr2.set("val", "E0E0E0")

            lin = ET.SubElement(grad_fill, _a("lin"))
            lin.set("ang", _direction_to_ang(gradient_direction))
            lin.set("scaled", "1")
        except Exception:
            return
    
    def _set_stroke_color_xml(self, shape, stroke_color: RGBColor):
        """Set stroke color via XML"""
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            ln_element = shape_element.find('.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
                if sp_pr is not None:
                    ln_element = ET.SubElement(sp_pr, _a('ln'))
                else:
                    return
            
            solid_fill = ln_element.find('.//a:solidFill', namespaces=NSMAP_DRAWINGML)
            if solid_fill is None:
                no_fill = ln_element.find('.//a:noFill', namespaces=NSMAP_DRAWINGML)
                if no_fill is not None:
                    ln_element.remove(no_fill)
                solid_fill = ET.SubElement(ln_element, _a('solidFill'))
            
            for color_elem in solid_fill.findall('.//a:srgbClr', namespaces=NSMAP_DRAWINGML):
                solid_fill.remove(color_elem)
            
            srgb = ET.SubElement(solid_fill, _a('srgbClr'))
            val = f"{stroke_color[0]:02X}{stroke_color[1]:02X}{stroke_color[2]:02X}"
            srgb.set('val', val)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set stroke color XML: {e}")
    
    def _set_edge_stroke_color_xml(self, shape, stroke_color: RGBColor):
        """Set edge stroke color via XML"""
        self._set_stroke_color_xml(shape, stroke_color)
    
    def _set_dash_pattern_xml(self, shape, dash_pattern: Optional[str]):
        """Set dash pattern via XML
        
        Args:
            shape: PowerPoint shape object
            dash_pattern: draw.io dash pattern name (e.g., "dashed", "dotted", "dashDot")
        """
        try:
            if not hasattr(shape, '_element'):
                return
            
            if not dash_pattern:
                return
            
            # Map draw.io dash pattern to PowerPoint prstDash value
            prst_dash = map_dash_pattern(dash_pattern)
            if not prst_dash or prst_dash == "solid":
                return
            
            shape_element = shape._element
            ln_element = shape_element.find(f'.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
                if sp_pr is not None:
                    ln_element = ET.SubElement(sp_pr, _a('ln'))
                else:
                    return
            
            # Remove existing prstDash
            for prst_dash_elem in ln_element.findall('.//a:prstDash', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(prst_dash_elem)
            
            # Add new prstDash
            prst_dash_elem = ET.SubElement(ln_element, _a('prstDash'))
            prst_dash_elem.set('val', prst_dash)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set dash pattern XML: {e}")
    
    def _set_arrow_heads_xml(
        self,
        shape,
        start_arrow: Optional[str],
        end_arrow: Optional[str],
        start_fill: bool,
        end_fill: bool,
        stroke_color: Optional[RGBColor],
        start_size_px: Optional[float],
        end_size_px: Optional[float],
    ):
        """Set arrows via XML

        Notes:
            DrawingML line-end elements (a:headEnd / a:tailEnd) are empty elements with attributes.
            They must NOT contain fill children (a:solidFill / a:noFill). PowerPoint uses the line
            formatting for arrow color. For "open" arrows, map to the appropriate line-end type.
        """
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            ln_element = shape_element.find(f'.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
                if sp_pr is not None:
                    ln_element = ET.SubElement(sp_pr, _a('ln'))
                else:
                    return
            
            # Remove existing arrows
            for head_end in ln_element.findall('.//a:headEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(head_end)
            for tail_end in ln_element.findall('.//a:tailEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(tail_end)
            
            # headEnd = line beginning (OOXML). draw.io startArrow
            if start_arrow:
                arrow_info = map_arrow_type_with_size(start_arrow, start_size_px)
                if arrow_info:
                    arrow_type, arrow_w, arrow_len = arrow_info
                    head_end = ET.SubElement(ln_element, _a('headEnd'))
                    head_end.set('type', arrow_type)
                    head_end.set('w', arrow_w)
                    head_end.set('len', arrow_len)
            
            # tailEnd = line end (OOXML). draw.io endArrow (arrow tip)
            if end_arrow:
                arrow_info = map_arrow_type_with_size(end_arrow, end_size_px)
                if arrow_info:
                    arrow_type, arrow_w, arrow_len = arrow_info
                    tail_end = ET.SubElement(ln_element, _a('tailEnd'))
                    tail_end.set('type', arrow_type)
                    tail_end.set('w', arrow_w)
                    tail_end.set('len', arrow_len)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set arrow XML: {e}")

    def _remove_arrowheads_xml(self, shape):
        """Remove arrowheads from a connector line via XML."""
        try:
            if not hasattr(shape, '_element'):
                return

            shape_element = shape._element
            ln_element = shape_element.find('.//a:ln', namespaces=NSMAP_DRAWINGML)
            if ln_element is None:
                return

            for head_end in ln_element.findall('.//a:headEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(head_end)
            for tail_end in ln_element.findall('.//a:tailEnd', namespaces=NSMAP_DRAWINGML):
                ln_element.remove(tail_end)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to remove arrowheads XML: {e}")

    def _remove_effect_ref_xml(self, shape):
        """Remove effectRef from connector style to avoid theme shadow."""
        try:
            if not hasattr(shape, '_element'):
                return
            shape_element = shape._element
            for effect_ref in list(shape_element.iter()):
                if effect_ref.tag.endswith('effectRef'):
                    parent = effect_ref.getparent()
                    if parent is not None:
                        parent.remove(effect_ref)
            for style in shape_element.findall('.//p:style', namespaces=NSMAP_BOTH):
                parent = style.getparent()
                if parent is not None:
                    parent.remove(style)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to remove effectRef XML: {e}")
    
    def _disable_shadow_xml(self, shape):
        """Disable shadow via XML (similar to legacy _disable_shadow_xml)"""
        try:
            if not hasattr(shape, '_element'):
                return
            
            shape_element = shape._element
            
            # Find spPr (shape properties) element
            sp_pr = shape_element.find('.//a:spPr', namespaces=NSMAP_DRAWINGML)
            if sp_pr is None:
                return
            
            # Remove existing effectLst (effects list) which contains shadow
            effect_lst = sp_pr.find('.//a:effectLst', namespaces=NSMAP_DRAWINGML)
            if effect_lst is not None:
                sp_pr.remove(effect_lst)
            
            # Also remove any shadow elements
            for shadow_elem in sp_pr.findall('.//a:outerShdw', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(shadow_elem)
            for shadow_elem in sp_pr.findall('.//a:innerShdw', namespaces=NSMAP_DRAWINGML):
                sp_pr.remove(shadow_elem)
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to disable shadow XML: {e}")

    def _set_cube_3d_rotation_xml(self, shape):
        """Apply 3D scene (camera + lightRig) and ratio adjustment to cube shape.
        - scene3d: isometricLeftDown + threePt light so the box looks tilted.
        - prstGeom adj=51666: makes the top face closer to a square (from reference timeline4_copy.pptx).
        """
        try:
            if not hasattr(shape, '_element'):
                return
            shape_element = shape._element
            # spPr is a direct child of p:sp with presentation namespace
            sp_pr = shape_element.find('.//p:spPr', namespaces=NSMAP_PRESENTATIONML)
            if sp_pr is None:
                for child in shape_element:
                    if child.tag.endswith('}spPr') or 'spPr' in child.tag:
                        sp_pr = child
                        break
            if sp_pr is None:
                return

            # Cube adjustment: adj=51666 so the top face is more square (reference file)
            prst_geom = sp_pr.find('.//a:prstGeom', namespaces=NSMAP_DRAWINGML)
            if prst_geom is not None and prst_geom.get('prst') == 'cube':
                av_lst = prst_geom.find('a:avLst', namespaces=NSMAP_DRAWINGML)
                if av_lst is None:
                    av_lst = ET.SubElement(prst_geom, _a('avLst'))
                # Remove existing adj if any, then add the reference value
                for gd in av_lst.findall('a:gd', namespaces=NSMAP_DRAWINGML):
                    if gd.get('name') == 'adj':
                        av_lst.remove(gd)
                        break
                gd_adj = ET.SubElement(av_lst, _a('gd'))
                gd_adj.set('name', 'adj')
                gd_adj.set('fmla', 'val 51666')

            if sp_pr.find('.//a:scene3d', namespaces=NSMAP_DRAWINGML) is not None:
                return
            scene3d = ET.SubElement(sp_pr, _a('scene3d'))
            camera = ET.SubElement(scene3d, _a('camera'))
            camera.set('prst', 'isometricLeftDown')
            light_rig = ET.SubElement(scene3d, _a('lightRig'))
            light_rig.set('rig', 'threePt')
            light_rig.set('dir', 't')
        except Exception as e:
            if self.logger:
                self.logger.debug(f"Failed to set cube 3D rotation XML: {e}")

    def _add_shape_image(self, slide, shape: ShapeElement):
        """
        Add image as a separate picture shape (not as shape fill)
        
        Args:
            slide: PowerPoint slide
            shape: ShapeElement with image data
        """
        if not shape.image:
            if self.logger:
                self.logger.debug("No image data in shape")
            return
        
        image_data = shape.image
        if self.logger:
            self.logger.debug(
                f"Processing image: data_uri={image_data.data_uri is not None}, file_path={image_data.file_path}"
            )

        left, top, width, height = self._compute_shape_geometry(shape)
        target_width_px = int(width / 9525) if width else None
        target_height_px = int(height / 9525) if height else None
        aws_icon_color_hex = None
        try:
            from ..stencil.aws_icons import is_aws_shape_type

            if is_aws_shape_type(shape.shape_type):
                fill = getattr(shape.style, "fill", None)
                if isinstance(fill, RGBColor):
                    aws_icon_color_hex = f"{fill[0]:02X}{fill[1]:02X}{fill[2]:02X}"
        except Exception:
            aws_icon_color_hex = None

        image_bytes, img_width_px, img_height_px, is_svg = prepare_image_for_pptx(
            data_uri=image_data.data_uri,
            file_path=image_data.file_path,
            shape_type=shape.shape_type,
            target_width_px=target_width_px,
            target_height_px=target_height_px,
            base_dpi=self.config.dpi if hasattr(self.config, "dpi") else 192.0,
            aws_icon_color_hex=aws_icon_color_hex,
        )

        if not image_bytes:
            if self.logger:
                self.logger.warning("No image bytes available after preparation")
            return

        if is_svg and self.logger and not getattr(self, "_svg_backend_logged", False):
            backend = getattr(self.config, "svg_backend", "cairosvg")
            self.logger.info(f"SVG to PNG: using {backend}")
            self._svg_backend_logged = True
        
        # Add image as a separate picture shape
        try:
            if img_width_px and img_height_px and img_width_px > 0 and img_height_px > 0:
                img_aspect = img_width_px / img_height_px
                original_height = height
                new_height = int(width / img_aspect)
                if new_height > 0:
                    height = new_height
                    if height != original_height:
                        top += int((original_height - height) / 2)
            
            # For image shapes, text is always placed below the image (y-axis direction)
            # The image uses the full shape area, and text will be added separately below
            
            # Create BytesIO object and ensure it's at the beginning
            image_stream = io.BytesIO(image_bytes)
            image_stream.seek(0)  # Reset stream position to beginning (required for add_picture)
            
            picture = slide.shapes.add_picture(
                image_stream,
                left, top, width, height
            )
            self._set_shape_name(picture, f"drawio2pptx:shape-image:{shape.id}" if shape.id else None)
            
            # Remove border from picture (no outline around image)
            try:
                picture.line.fill.background()
            except Exception:
                pass
            
            # Apply transform if needed
            if shape.transform.rotation:
                picture.rotation = shape.transform.rotation
            
            if self.logger:
                self.logger.debug(f"Successfully added image as picture shape at ({left}, {top}), size ({width}, {height})")
        except Exception as e:
            if self.logger:
                self.logger.warning(f"Failed to add image as picture shape: {e}")
                import traceback
                self.logger.debug(traceback.format_exc())
            # Also print to stderr for debugging
            import sys
            print(f"ERROR: Failed to add image: {e}", file=sys.stderr)
            import traceback
            traceback.print_exc(file=sys.stderr)
    
